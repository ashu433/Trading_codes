import intraday_initial_final_sentiment
import matplotlib.pyplot as plt
import collections 
import collections.abc
import pptx
from pptx import Presentation
import pandas as pd
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN



def Every_day_initial_final_Oi_buildup(number_of_rows,closing_price,Date):

    path="D:/ashu/Finance/algo_trading/Option_chain_data/"
    file=path+"full_and_finial_file.xlsx"
    number_of_rows=number_of_rows
    closing_price=closing_price
    opening_price=18294.8

    call,put=intraday_initial_final_sentiment.over_all_OI_build_up_throughout(file,closing_price)

    call.to_csv(path+"Call_OI_build_up.csv",index=False)
    put.to_csv(path+"Put_OI_build_up.csv",index=False)

    # intraday_initial_final_sentiment.start_of_the_day_sentiment_call(file,77,4,ATM_start,ATM_end)


    # intraday_initial_final_sentiment.start_of_the_day_sentiment_put(file,77,4,ATM_start,ATM_end)

    intraday_initial_final_sentiment.final_45_minutes_start_call(file,number_of_rows,closing_price)

    intraday_initial_final_sentiment.final_45_minutes_start_put(file,number_of_rows,closing_price)



    ############################### PPT Making #############################

    prs=Presentation()

    ############################################ Slide 3 #################################################

    df_1=pd.read_csv(path+"Call_OI_build_up.csv")
    df_1["Change_in_Premium"] = df_1["Change_in_Premium"].round(2)
    new_data_frame=df_1
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    title_3=slide3.shapes.title
    title_3.text=f"EOD Call Side Conclusion {Date}"
    # create a table on the slide
    rows = new_data_frame.shape[0] + 1  # add 1 for header row
    cols = new_data_frame.shape[1]
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(8.3)
    height = Inches(2)
    table = slide3.shapes.add_table(rows, cols, left, top, width, height).table

    # set the header row
    header_row = table.rows[0]
    for i, col_name in enumerate(new_data_frame.columns):
        header_row.cells[i].text = col_name
        header_row.cells[i].text_frame.paragraphs[0].font.size = pptx.util.Pt(12)

    # set the data rows
    for i in range(new_data_frame.shape[0]):
        row = table.rows[i+1]  # add 1 for header row
        for j in range(new_data_frame.shape[1]):
            cell = row.cells[j]
            cell.text = str(new_data_frame.iloc[i, j])
            cell.text_frame.paragraphs[0].font.size = pptx.util.Pt(10)  # set font size for data cells



    ############################################ Slide 4 #################################################

    df_2=pd.read_csv(path+"Put_OI_build_up.csv")
    df_2["Change_in_Premium"] = df_2["Change_in_Premium"].round(2)
    new_data_frame=df_2
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    title_4=slide4.shapes.title
    title_4.text=f"EOD Put Side Conclusion {Date}"
    # create a table on the slide
    rows = new_data_frame.shape[0] + 1  # add 1 for header row
    cols = new_data_frame.shape[1]
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(8.3)
    height = Inches(2)
    table = slide4.shapes.add_table(rows, cols, left, top, width, height).table

    # set the header row
    header_row = table.rows[0]
    for i, col_name in enumerate(new_data_frame.columns):
        header_row.cells[i].text = col_name
        header_row.cells[i].text_frame.paragraphs[0].font.size = pptx.util.Pt(12)

    # set the data rows
    for i in range(new_data_frame.shape[0]):
        row = table.rows[i+1]  # add 1 for header row
        for j in range(new_data_frame.shape[1]):
            cell = row.cells[j]
            cell.text = str(new_data_frame.iloc[i, j])
            cell.text_frame.paragraphs[0].font.size = pptx.util.Pt(10)  # set font size for data cells



    ################################### Slide 5 ########################################################

    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    title_5=slide5.shapes.title
    title_5.text=f"Final 45 mins Call side {Date}"

    image_file1 = path+'Maximum_OI_on_the_Call_side_VS_strike_price.png'
    #image_title1 = 'Micro and Macro sentiment'

    # Add the first image to the slide
    left1 = Inches(0)
    top1 = Inches(1.5)
    height1 = Inches(5)
    picture1 = slide5.shapes.add_picture(image_file1, left1, top1, height=height1)

    image_file1 = path+'Premium_on_Call_side_VS_strike_price.png'
    #image_title1 = 'Micro and Macro sentiment'

    # Add the first image to the slide
    left1 = Inches(5)
    top1 = Inches(1.5)
    height1 = Inches(5)
    picture1 = slide5.shapes.add_picture(image_file1, left1, top1, height=height1)


    ################################### Slide 6 ########################################################

    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    title_6=slide6.shapes.title
    title_6.text=f"Final 45 mins Put side {Date}"

    image_file1 = path+'Maximum_OI_on_the_Put_side_VS_strike_price.png'
    #image_title1 = 'Micro and Macro sentiment'

    # Add the first image to the slide
    left1 = Inches(0)
    top1 = Inches(1.5)
    height1 = Inches(5)
    picture1 = slide6.shapes.add_picture(image_file1, left1, top1, height=height1)

    image_file1 = path+'Premium_on_Put_side_VS_strike_price.png'
    #image_title1 = 'Micro and Macro sentiment'

    # Add the first image to the slide
    left1 = Inches(5)
    top1 = Inches(1.5)
    height1 = Inches(5)
    picture1 = slide6.shapes.add_picture(image_file1, left1, top1, height=height1)

    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    title_7=slide7.shapes.title
    title_7.text=f"Over All conclusion {Date}"

    bullet_points=slide7.shapes
    points=bullet_points.placeholders[1]
    points.text = 'Overall Sentiment:\nFinal 45 mins Sentiment:\nMy Vote:\nResult:'


    prs.save(path+"Overall_OI_data_conclusion"".pptx")


if __name__=='__main__':  
    pass