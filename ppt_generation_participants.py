import collections 
import collections.abc
import pptx
from pptx import Presentation
import pandas as pd
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


def ppt_generation(Date):
    prs=Presentation()
    path="D:/ashu/Finance/algo_trading/images/"

    Date=Date

    ############################################ Title_slide slide 1 #################################################

    slide1=prs.slides.add_slide(prs.slide_layouts[0])
    title=slide1.shapes.title
    title.text="Daily market Report of different participants"

    subtitle = slide1.placeholders[1]
    subtitle.text = Date

    ############################################ Slide 2 #################################################

    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    title_2=slide2.shapes.title
    title_2.text="Content"

    bullet_points=slide2.shapes
    points=bullet_points.placeholders[1]
    points.text = 'FII Daily market Report\nParticipants wise daily market report in Futures\nParticipants wise daily market report in Options\nOver all market report'

    ############################################ Slide 3 #################################################

    df_1=pd.read_csv(path+"FII_stats.csv")
    new_data_frame=df_1.iloc[:11]

    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    title_3=slide3.shapes.title
    title_3.text="Fii Stats Report in Futures and Options"
    bullet_points_3=slide3.shapes
    points_3=bullet_points_3.placeholders[1]
    points_3.text="FII F&O stats in the index"
    # create a table on the slide
    rows = new_data_frame.shape[0] + 1  # add 1 for header row
    cols = new_data_frame.shape[1]
    left = Inches(1)
    top = Inches(3)
    width = Inches(8)
    height = Inches(2)
    table = slide3.shapes.add_table(rows, cols, left, top, width, height).table

    # set the header row
    header_row = table.rows[0]
    for i, col_name in enumerate(new_data_frame.columns):
        header_row.cells[i].text = col_name
        header_row.cells[i].text_frame.paragraphs[0].font.size = pptx.util.Pt(12)

    # set the data rows
    for i in range(new_data_frame.shape[0]):
        row = table.rows[i+1]  # add 1 for header row
        for j in range(new_data_frame.shape[1]):
            cell = row.cells[j]
            cell.text = str(new_data_frame.iloc[i, j])
            cell.text_frame.paragraphs[0].font.size = pptx.util.Pt(10)  # set font size for data cells


    ################################### Slide 4 ########################################################

    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    title_4=slide4.shapes.title
    title_4.text="Fii Index share in F&O and Macro and Micro view"

    image_file1 = path+'Micro and Macro sentiment of the FII.png'
    #image_title1 = 'Micro and Macro sentiment'

    # Add the first image to the slide
    left1 = Inches(0)
    top1 = Inches(2.5)
    height1 = Inches(3.8)
    picture1 = slide4.shapes.add_picture(image_file1, left1, top1, height=height1)

    image_file1 = path+'FII Index share in each instrument.png'
    #image_title1 = 'Micro and Macro sentiment'

    # Add the first image to the slide
    left1 = Inches(5)
    top1 = Inches(2.5)
    height1 = Inches(3.75)
    picture1 = slide4.shapes.add_picture(image_file1, left1, top1, height=height1)


    ################################### Slide 5 ########################################################

    new_data_frame_fut=pd.read_csv(path+"Participant_Futures.csv")

    new_data_frame_fut["percentage_long"]=new_data_frame_fut["percentage_long"].round(2)
    new_data_frame_fut["percentage_short"]=new_data_frame_fut["percentage_short"].round(2)
    new_data_frame_fut["Long_vs_Short"]=new_data_frame_fut["Long_vs_Short"].round(2)

    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    title_5=slide5.shapes.title
    title_5.text="Participants report in Futures"
    bullet_points_5=slide5.shapes
    points_5=bullet_points_5.placeholders[1]
    points_5.text="Participant wise Open Interest in Equity Derivatives as on "+Date
    # create a table on the slide
    rows = new_data_frame_fut.shape[0] + 1  # add 1 for header row
    cols = new_data_frame_fut.shape[1]
    left = Inches(0.5)
    top = Inches(3)
    width = Inches(9)
    height = Inches(2)
    table = slide5.shapes.add_table(rows, cols, left, top, width, height).table

    # set the header row
    header_row = table.rows[0]
    for i, col_name in enumerate(new_data_frame_fut.columns):
        header_row.cells[i].text = col_name
        header_row.cells[i].text_frame.paragraphs[0].font.size = pptx.util.Pt(12)

    # set the data rows
    for i in range(new_data_frame_fut.shape[0]):
        row = table.rows[i+1]  # add 1 for header row
        for j in range(new_data_frame_fut.shape[1]):
            cell = row.cells[j]
            cell.text = str(new_data_frame_fut.iloc[i, j])
            cell.text_frame.paragraphs[0].font.size = pptx.util.Pt(10)  # set font size for data cells


    ################################### Slide 6 ########################################################

    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    title_6=slide6.shapes.title
    title_6.text="% OI of the different participants in Futures"

    image_file2 = path+'Total_Open_Interest_of_the_different_participants_Futures.png'
    #image_title1 = 'Micro and Macro sentiment'

    # Add the first image to the slide
    left1 = Inches(1.2)
    top1 = Inches(2)
    height1 = Inches(5.2)
    picture1 = slide6.shapes.add_picture(image_file2, left1, top1, height=height1)


    ################################### Slide 7 ########################################################

    slide7 = prs.slides.add_slide(prs.slide_layouts[5])
    title_7=slide7.shapes.title
    title_7.text="Total Long pogitions of different clients in Futures"

    image_file3 = path+'Total_Long_Pogition_in_Futures.png'
    #image_title1 = 'Micro and Macro sentiment'

    # Add the first image to the slide
    left1 = Inches(0)
    top1 = Inches(2.5)
    height1 = Inches(3.8)
    picture1 = slide7.shapes.add_picture(image_file3, left1, top1, height=height1)

    image_file3 = path+'Total_Long_Pogition_of_big_clients_Futures.png'
    #image_title1 = 'Micro and Macro sentiment'

    # Add the first image to the slide
    left1 = Inches(5)
    top1 = Inches(2.5)
    height1 = Inches(3.75)
    picture1 = slide7.shapes.add_picture(image_file3, left1, top1, height=height1)



    ################################### Slide 8 ########################################################

    slide8 = prs.slides.add_slide(prs.slide_layouts[5])
    title_8=slide8.shapes.title
    title_8.text="Total Shot pogitions of different clients in Futures"

    image_file4 = path+'Total_Short_Pogition_in_Futures.png'
    #image_title1 = 'Micro and Macro sentiment'

    # Add the first image to the slide
    left1 = Inches(0)
    top1 = Inches(2.5)
    height1 = Inches(3.8)
    picture1 = slide8.shapes.add_picture(image_file4, left1, top1, height=height1)

    image_file4 = path+'Total_Shot_Pogition_of_big_clients_Futures.png'
    #image_title1 = 'Micro and Macro sentiment'

    # Add the first image to the slide
    left1 = Inches(5)
    top1 = Inches(2.5)
    height1 = Inches(3.75)
    picture1 = slide8.shapes.add_picture(image_file4, left1, top1, height=height1)

    ################################### Slide 9 ########################################################

    new_data_frame_opt=pd.read_csv(path+"Participant_Option.csv")

    new_data_frame_opt["percentage_long_option"]=new_data_frame_opt["percentage_long_option"].round(2)
    new_data_frame_opt["percentage_short_option"]=new_data_frame_opt["percentage_short_option"].round(2)
    new_data_frame_opt["Long_vs_Short"]=new_data_frame_opt["Long_vs_Short"].round(2)

    slide9 = prs.slides.add_slide(prs.slide_layouts[1])
    title_9=slide9.shapes.title
    title_9.text="Participants report in Options"
    bullet_points_9=slide9.shapes
    points_9=bullet_points_9.placeholders[1]
    points_9.text="Participant wise Open Interest in Equity Derivatives as on "+Date
    # create a table on the slide
    rows = new_data_frame_opt.shape[0] + 1  # add 1 for header row
    cols = new_data_frame_opt.shape[1]
    left = Inches(0)
    top = Inches(3)
    width = Inches(10)
    height = Inches(2)
    table = slide9.shapes.add_table(rows, cols, left, top, width, height).table

    # set the header row
    header_row = table.rows[0]
    for i, col_name in enumerate(new_data_frame_opt.columns):
        header_row.cells[i].text = col_name
        header_row.cells[i].text_frame.paragraphs[0].font.size = pptx.util.Pt(12)

    # set the data rows
    for i in range(new_data_frame_opt.shape[0]):
        row = table.rows[i+1]  # add 1 for header row
        for j in range(new_data_frame_opt.shape[1]):
            cell = row.cells[j]
            cell.text = str(new_data_frame_opt.iloc[i, j])
            cell.text_frame.paragraphs[0].font.size = pptx.util.Pt(10)  # set font size for data cells

    ################################### Slide 10 ########################################################

    slide10 = prs.slides.add_slide(prs.slide_layouts[5])
    title_10=slide10.shapes.title
    title_10.text="% OI of the different participants in Options"

    image_file5 = path+'Total_Open_Interest_of_the_different_participants_Options.png'
    #image_title1 = 'Micro and Macro sentiment'

    # Add the first image to the slide
    left1 = Inches(1.2)
    top1 = Inches(2)
    height1 = Inches(5.2)
    picture1 = slide10.shapes.add_picture(image_file5, left1, top1, height=height1)


    ################################### Slide 11 ########################################################

    slide11 = prs.slides.add_slide(prs.slide_layouts[5])
    title_11=slide11.shapes.title
    title_11.text="Total Long pogitions of different clients in Options"

    image_file6 = path+'Total_Long_pogition_of_different_participants_Options.png'
    #image_title1 = 'Micro and Macro sentiment'

    # Add the first image to the slide
    left1 = Inches(0)
    top1 = Inches(2.5)
    height1 = Inches(3.8)
    picture1 = slide11.shapes.add_picture(image_file6, left1, top1, height=height1)

    image_file6 = path+'Total_Long_Pogition_of_big_clients_Options.png'
    #image_title1 = 'Micro and Macro sentiment'

    # Add the first image to the slide
    left1 = Inches(5)
    top1 = Inches(2.5)
    height1 = Inches(3.75)
    picture1 = slide11.shapes.add_picture(image_file6, left1, top1, height=height1)

    ################################### Slide 12 ########################################################

    slide12 = prs.slides.add_slide(prs.slide_layouts[5])
    title_12=slide12.shapes.title
    title_12.text="Total Shot pogitions of different clients in Options"

    image_file7 = path+'Total_Short_Pogition_in_Options.png'
    #image_title1 = 'Micro and Macro sentiment'

    # Add the first image to the slide
    left1 = Inches(0)
    top1 = Inches(2.5)
    height1 = Inches(3.8)
    picture1 = slide12.shapes.add_picture(image_file7, left1, top1, height=height1)

    image_file7 = path+'Total_Shot_Pogition_of_big_clients_Options.png'
    #image_title1 = 'Micro and Macro sentiment'

    # Add the first image to the slide
    left1 = Inches(5)
    top1 = Inches(2.5)
    height1 = Inches(3.75)
    picture1 = slide12.shapes.add_picture(image_file7, left1, top1, height=height1)

    ################################### Slide 13 ########################################################

    slide13 = prs.slides.add_slide(prs.slide_layouts[5])
    title_13=slide13.shapes.title
    title_13.text="Over all % Long and the Short pogition of the different client in F&O"

    image_file8 = path+'%_long_and_short_of_the_each_Participants_OI_in_Futures.png'
    #image_title1 = 'Micro and Macro sentiment'

    # Add the first image to the slide
    left1 = Inches(0)
    top1 = Inches(2.5)
    height1 = Inches(5)
    picture1 = slide13.shapes.add_picture(image_file8, left1, top1, height=height1)

    image_file8 = path+'%_long_and_short_of_the_each_Participants_OI_in_Options.png'
    #image_title1 = 'Micro and Macro sentiment'

    # Add the first image to the slide
    left1 = Inches(5)
    top1 = Inches(2.5)
    height1 = Inches(5)
    picture1 = slide13.shapes.add_picture(image_file8, left1, top1, height=height1)

    prs.save(path+"Daily_market_Report_"+Date+".pptx")



if __name__=='__main__':  
    pass